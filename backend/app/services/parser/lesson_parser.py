
import os
import uuid
import requests
import tempfile
import base64
import json
import time
import asyncio
import win32com.client
import pythoncom
from typing import Dict, Any, List, Optional, Tuple
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader
from langchain_core.messages import HumanMessage
from backend.app.utils.llm_pool import (
    LLMPoolConfig,
    initialize_pool,
    get_llm_client,
    release_llm_client,
)
from backend.app.core.config import settings
from backend.app.schemas.parser import Chapter, SubChapter, FileInfo, StructurePreview, PageElement, BoundingBox

# Import pdfminer.six components
from pdfminer.high_level import extract_pages
from pdfminer.layout import LTTextContainer, LTChar, LTTextLine, LTTextBox

class LessonParserService:
    def __init__(
        self,
        vision_api_key: Optional[str] = None,
        vision_base_url: Optional[str] = None,
        vision_model: Optional[str] = None,
    ):
        self.vision_api_key = (
            vision_api_key
            or settings.EFFECTIVE_IMAGE_VLM_API_KEY
            or settings.EFFECTIVE_QWEN_API_KEY
            or os.getenv("OPENAI_API_KEY")
            or ""
        )
        self.vision_base_url = (
            vision_base_url
            or settings.EFFECTIVE_IMAGE_VLM_BASE_URL
            or settings.EFFECTIVE_QWEN_BASE_URL
            or "https://dashscope.aliyuncs.com/compatible-mode/v1"
        )
        self.vision_model = (
            vision_model
            or settings.EFFECTIVE_IMAGE_VLM_MODEL
            or settings.QWEN_VISION_MODEL
            or "qwen-vl-plus"
        )
        self.image_parse_concurrency = int(os.getenv("IMAGE_PARSE_CONCURRENCY", "8"))
        self.image_pool_max_connections = int(os.getenv("IMAGE_POOL_MAX_CONNECTIONS", "100"))
        self.image_pool_min_idle = int(os.getenv("IMAGE_POOL_MIN_IDLE", "10"))
        self.progress_enabled = os.getenv("PARSER_PROGRESS", "1") == "1"
        self.vision_prompt = (
            "# Role\n"
            "你是一位资深的PPT内容分析与视觉理解专家。你具备敏锐的洞察力，能够精准识别幻灯片中的视觉元素，并深刻理解图片与文字之间的逻辑联系。\n\n"
            "# Task\n"
            "请分析提供的slides与图片字典json文件，完成以下任务：\n"
            "给每一个图片字典中的图片信息relationship与content_data信息\n\n"
            "# Constraints\n"
            "- 必须仅输出标准的JSON格式字符串。\n"
            "- 不要输出任何Markdown标记（如 ```json ... ```）、开场白或解释性文字。\n"
            "- 确保JSON可以被代码直接解析（无语法错误）。\n"
            "- relationship不要超过10个字。\n"
            "- content_data不要超过20个字。\n\n"
            "# Output Format\n"
            "{\n"
            '    "relationship": "在此处描述图片与文字的逻辑关系（例如：图片展示了文字中提到的数据增长趋势，作为直观的证据支撑...）",\n'
            '    "content_data": "在此处详细描述图片的具体视觉内容（例如：图片包含三个主要部分，左侧是...，右侧是...，颜色主要为...）"\n'
            "}\n"
        )

    def _clip_text(self, text: str, max_len: int) -> str:
        return (text or "").strip()[:max_len]

    def _print_progress(self, stage: str, current: int, total: int) -> None:
        if not self.progress_enabled:
            return
        total = max(1, total)
        percent = int((current / total) * 100)
        print(f"\r[{stage}] {current}/{total} ({percent}%)", end="", flush=True)
        if current >= total:
            print("", flush=True)

    def _parse_image_json_output(self, raw_text: str) -> Dict[str, str]:
        fallback = {"relationship": "辅助说明", "content_data": "图像解析失败"}
        if not raw_text:
            return fallback

        cleaned = raw_text.strip()
        try:
            payload = json.loads(cleaned)
        except json.JSONDecodeError:
            start = cleaned.find("{")
            end = cleaned.rfind("}")
            if start == -1 or end == -1 or end <= start:
                return fallback
            try:
                payload = json.loads(cleaned[start : end + 1])
            except json.JSONDecodeError:
                return fallback

        relationship = self._clip_text(str(payload.get("relationship", "")), 10)
        content_data = self._clip_text(str(payload.get("content_data", "")), 20)
        if not relationship:
            relationship = "辅助说明"
        if not content_data:
            content_data = "图像内容不明"
        return {"relationship": relationship, "content_data": content_data}

    def _image_to_data_url(self, image_blob: bytes, ext: str) -> str:
        safe_ext = (ext or "png").lower().lstrip(".")
        if safe_ext == "jpg":
            safe_ext = "jpeg"
        b64 = base64.b64encode(image_blob).decode("ascii")
        return f"data:image/{safe_ext};base64,{b64}"

    async def _describe_one_image_async(self, image_blob: bytes, ext: str, sem: asyncio.Semaphore) -> Tuple[Dict[str, str], float]:
        client = None
        start_ts = time.perf_counter()
        try:
            async with sem:
                client = get_llm_client(model=self.vision_model)
                message = HumanMessage(
                    content=[
                        {"type": "text", "text": self.vision_prompt},
                        {"type": "image_url", "image_url": {"url": self._image_to_data_url(image_blob, ext)}},
                    ]
                )
                response = await client.ainvoke([message])
                content = response.content
                if isinstance(content, list):
                    text_parts = [c.get("text", "") for c in content if isinstance(c, dict)]
                    content = "\n".join([p for p in text_parts if p])
                return self._parse_image_json_output(str(content or "")), (time.perf_counter() - start_ts)
        except Exception:
            return {"relationship": "解析异常", "content_data": "图片理解失败"}, (time.perf_counter() - start_ts)
        finally:
            if client is not None:
                try:
                    release_llm_client(client)
                except Exception:
                    pass

    def _describe_images_parallel(self, image_inputs: List[Tuple[bytes, str]]) -> List[Dict[str, str]]:
        if not image_inputs:
            return []

        if not self.vision_api_key:
            return [{"relationship": "辅助说明", "content_data": "未配置图片模型"} for _ in image_inputs]

        # Inject runtime vision config so llm_pool can use the provided API key/model.
        settings.QWEN_API_KEY = self.vision_api_key
        settings.QWEN_BASE_URL = self.vision_base_url
        settings.QWEN_VISION_MODEL = self.vision_model

        pool_start_ts = time.perf_counter()
        try:
            initialize_pool(
                LLMPoolConfig(
                    max_connections=max(1, self.image_pool_max_connections),
                    min_idle_connections=max(1, min(self.image_pool_min_idle, self.image_pool_max_connections)),
                    connection_timeout=120.0,
                    idle_timeout=300.0,
                    health_check_interval=60.0,
                )
            )
        except Exception:
            return [{"relationship": "辅助说明", "content_data": "未配置图片模型"} for _ in image_inputs]
        pool_elapsed = time.perf_counter() - pool_start_ts
        if self.progress_enabled:
            print(f"[计时] llm_pool建立耗时: {pool_elapsed:.3f}s")

        async def _runner() -> List[Dict[str, str]]:
            first_parse_started_ts: Optional[float] = None
            sem = asyncio.Semaphore(max(1, self.image_parse_concurrency))
            async def _run_one(idx: int, blob: bytes, ext: str) -> Tuple[int, Dict[str, str]]:
                nonlocal first_parse_started_ts
                if first_parse_started_ts is None:
                    first_parse_started_ts = time.perf_counter()
                parsed, _ = await self._describe_one_image_async(blob, ext, sem)
                return idx, parsed

            tasks = [
                asyncio.create_task(_run_one(idx, blob, ext))
                for idx, (blob, ext) in enumerate(image_inputs)
            ]
            total = len(tasks)
            done = 0
            results: List[Optional[Dict[str, str]]] = [None] * total
            for task in asyncio.as_completed(tasks):
                idx, res = await task
                results[idx] = res
                done += 1
                self._print_progress("图片解析", done, total)
            if self.progress_enabled and total > 0:
                batch_end_ts = time.perf_counter()
                batch_elapsed = batch_end_ts - (first_parse_started_ts or batch_end_ts)
                print(
                    f"[计时] llm生成耗时: 批次总耗时 {batch_elapsed:.3f}s, "
                    f"单图折算 {batch_elapsed/total:.3f}s, 图片数 {total}"
                )
            return [r or {"relationship": "解析异常", "content_data": "图片理解失败"} for r in results]

        try:
            return asyncio.run(_runner())
        except RuntimeError:
            # In case an event loop is already running (rare in sync parser context).
            loop = asyncio.new_event_loop()
            try:
                return loop.run_until_complete(_runner())
            finally:
                loop.close()

    def parse_courseware(self, file_url: str, file_type: str) -> Dict[str, Any]:
        """
        Main entry point for parsing courseware.
        Returns a dictionary matching the ParseData schema structure (minus parseId/taskStatus).
        """
        start_ts = time.perf_counter()
        # 1. Download the file
        temp_file_path = self._download_file(file_url)
        converted_file_path = None
        
        try:
            parse_result: Dict[str, Any]
            if file_type.lower() in ['ppt', 'pptx']:
                # If it's a .ppt file, convert to .pptx first using COM
                if temp_file_path.lower().endswith('.ppt'):
                    converted_file_path = self._convert_ppt_to_pptx(temp_file_path)
                    parse_result = self._parse_ppt(converted_file_path, file_url)
                else:
                    parse_result = self._parse_ppt(temp_file_path, file_url)
            elif file_type.lower() == 'pdf':
                parse_result = self._parse_pdf(temp_file_path, file_url)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
            parse_result["parseElapsedSeconds"] = round(time.perf_counter() - start_ts, 3)
            if self.progress_enabled:
                print(f"[计时] 解析总耗时: {parse_result['parseElapsedSeconds']:.3f}s")
            return parse_result
        finally:
            # Cleanup temp file
            if os.path.exists(temp_file_path):
                # Only delete if it was downloaded (not local original)
                if temp_file_path != file_url:
                    try:
                        os.remove(temp_file_path)
                    except:
                        pass
            if converted_file_path and os.path.exists(converted_file_path):
                try:
                    os.remove(converted_file_path)
                except:
                    pass

    def _download_file(self, url: str) -> str:
        # Check if it's a local path first (for testing/sandbox)
        if os.path.exists(url):
            return url
            
        response = requests.get(url, stream=True)
        response.raise_for_status()
        
        # Create a temporary file
        suffix = os.path.splitext(url)[1]
        if not suffix:
             # Default fallback if no extension in URL
            suffix = ".tmp" 
            
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            for chunk in response.iter_content(chunk_size=8192):
                tmp.write(chunk)
            return tmp.name

    def _convert_ppt_to_pptx(self, ppt_path: str) -> str:
        """
        Convert .ppt to .pptx using Windows COM interface.
        """
        pythoncom.CoInitialize()
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        
        abs_ppt_path = os.path.abspath(ppt_path)
        output_path = os.path.splitext(abs_ppt_path)[0] + "_converted.pptx"
        
        try:
            # Use headless mode with caution, sometimes visible=1 works better
            # Open(FileName, ReadOnly, Untitled, WithWindow)
            deck = powerpoint.Presentations.Open(abs_ppt_path, True, False, False)
            deck.SaveAs(output_path, 24) # 24 is ppSaveAsOpenXMLPresentation
            deck.Close()
            return output_path
        except Exception as e:
            raise RuntimeError(f"Failed to convert PPT to PPTX: {str(e)}")
        finally:
            if powerpoint.Presentations.Count == 0:
                powerpoint.Quit()

    def _parse_ppt(self, file_path: str, original_url: str) -> Dict[str, Any]:
        prs = Presentation(file_path)
        file_size = os.path.getsize(file_path)
        page_count = len(prs.slides)
        file_name = os.path.basename(original_url)

        chapters = []
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        current_chapter = Chapter(
            chapterId=str(uuid.uuid4()),
            chapterName="默认章节",
            subChapters=[]
        )
        
        all_image_inputs: List[Tuple[bytes, str]] = []
        all_image_targets: List[Tuple[SubChapter, float, float, float, float]] = []

        for i, slide in enumerate(prs.slides):
            title = f"第 {i+1} 页"
            if slide.shapes.title:
                title = slide.shapes.title.text
            
            elements = []
            # Extract elements with coordinates
            for shape in slide.shapes:
                # Some PPT shapes (e.g., placeholders/groups) may not expose geometry.
                if any(getattr(shape, attr, None) is None for attr in ("left", "top", "width", "height")):
                    continue

                # Calculate relative coordinates
                # shape.left, shape.top are in EMUs
                # slide_width, slide_height are in EMUs
                x = shape.left / slide_width
                y = shape.top / slide_height
                w = shape.width / slide_width
                h = shape.height / slide_height

                if shape.has_text_frame and shape.text and shape.text.strip():
                    elements.append(PageElement(
                        type="text",
                        content=shape.text.strip(),
                        bbox=BoundingBox(x=x, y=y, width=w, height=h)
                    ))
                    continue

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        all_image_inputs.append((image.blob, image.ext or "png"))
                        all_image_targets.append((None, x, y, w, h))  # subchapter filled after creation
                    except Exception:
                        elements.append(PageElement(
                            type="image",
                            content="图片读取失败",
                            relationship="解析异常",
                            content_data="图片读取失败",
                            bbox=BoundingBox(x=x, y=y, width=w, height=h)
                        ))

            sub_chapter = SubChapter(
                subChapterId=str(uuid.uuid4()),
                subChapterName=title.strip() if title else f"Page {i+1}",
                isKeyPoint=False,
                pageRange=f"{i+1}",
                elements=elements
            )
            # Bind pending image targets created on this slide to current subchapter.
            for idx in range(len(all_image_targets) - 1, -1, -1):
                target_sub, x, y, w, h = all_image_targets[idx]
                if target_sub is None:
                    all_image_targets[idx] = (sub_chapter, x, y, w, h)
                else:
                    break
            current_chapter.subChapters.append(sub_chapter)
            self._print_progress("页面解析", i + 1, page_count)
            
        # Cross-page parallel image parsing in one batch.
        all_image_descs = self._describe_images_parallel(all_image_inputs)
        for (sub_chapter, x, y, w, h), image_desc in zip(all_image_targets, all_image_descs):
            if sub_chapter is None:
                continue
            sub_chapter.elements.append(PageElement(
                type="image",
                content=image_desc["content_data"],
                relationship=image_desc["relationship"],
                content_data=image_desc["content_data"],
                bbox=BoundingBox(x=x, y=y, width=w, height=h)
            ))

        chapters.append(current_chapter)

        return {
            "fileInfo": FileInfo(fileName=file_name, fileSize=file_size, pageCount=page_count),
            "structurePreview": StructurePreview(chapters=chapters)
        }

    def _parse_pdf(self, file_path: str, original_url: str) -> Dict[str, Any]:
        """
        Parse PDF using pdfminer.six to extract text blocks with coordinates.
        """
        # Get basic info
        file_size = os.path.getsize(file_path)
        file_name = os.path.basename(original_url)
        
        # Use pypdf for page count (fast)
        reader = PdfReader(file_path)
        page_count = len(reader.pages)
        
        chapters = []
        current_chapter = Chapter(
            chapterId=str(uuid.uuid4()),
            chapterName="PDF 内容",
            subChapters=[]
        )
        
        # Extract pages with layout analysis
        # This is slower but provides coordinates
        for i, page_layout in enumerate(extract_pages(file_path)):
            page_width = page_layout.width
            page_height = page_layout.height
            
            elements = []
            page_text_content = ""
            
            for element in page_layout:
                if isinstance(element, LTTextContainer):
                    text = element.get_text().strip()
                    if not text:
                        continue
                        
                    page_text_content += text + "\n"
                    
                    # Coordinate system in PDF is bottom-left origin
                    # We need top-left origin for frontend consistency (0.0 to 1.0)
                    
                    # element.x0, element.y0 (bottom-left)
                    # element.x1, element.y1 (top-right)
                    
                    # Calculate relative coordinates (Top-Left Origin)
                    # x = x0 / width
                    # y = (height - y1) / height  <-- Flip Y axis
                    # w = (x1 - x0) / width
                    # h = (y1 - y0) / height
                    
                    x = element.x0 / page_width
                    y = (page_height - element.y1) / page_height
                    w = (element.x1 - element.x0) / page_width
                    h = (element.y1 - element.y0) / page_height
                    
                    # Clamp values to 0-1 just in case
                    x = max(0.0, min(1.0, x))
                    y = max(0.0, min(1.0, y))
                    w = max(0.0, min(1.0, w))
                    h = max(0.0, min(1.0, h))
                    
                    elements.append(PageElement(
                        type="text",
                        content=text,
                        bbox=BoundingBox(x=x, y=y, width=w, height=h)
                    ))

            # Determine title from first line of text
            lines = page_text_content.split('\n')
            title = lines[0] if lines else f"第 {i+1} 页"
            if len(title) > 50:
                title = title[:50] + "..."

            sub_chapter = SubChapter(
                subChapterId=str(uuid.uuid4()),
                subChapterName=title.strip(),
                isKeyPoint=False,
                pageRange=f"{i+1}",
                elements=elements
            )
            current_chapter.subChapters.append(sub_chapter)
            self._print_progress("页面解析", i + 1, page_count)
            
        chapters.append(current_chapter)
        
        return {
            "fileInfo": FileInfo(fileName=file_name, fileSize=file_size, pageCount=page_count),
            "structurePreview": StructurePreview(chapters=chapters)
        }

parser_service = LessonParserService()
